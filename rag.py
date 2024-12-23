# Import necessary libraries
import streamlit as st
from langchain_community.document_loaders import PyPDFLoader, TextLoader  # Updated imports
from langchain_community.vectorstores import FAISS  # Updated imports
from langchain.schema import Document
from langchain.chains import RetrievalQA
from langchain_community.llms import OpenAI, Ollama  # Updated imports
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
from transformers import pipeline
from langchain_community.embeddings import HuggingFaceEmbeddings  # Updated imports
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain
from langchain_core.messages import HumanMessage,AIMessage
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.prompts import MessagesPlaceholder
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.chains import create_retrieval_chain
from langchain.chains import create_history_aware_retriever
from langchain.memory import ConversationBufferMemory
from langchain_community.llms import HuggingFacePipeline  # Updated imports
import os
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from transformers import BartForConditionalGeneration, BartTokenizer
from difflib import HtmlDiff, SequenceMatcher
import tempfile
from transformers import AutoModelForCausalLM, AutoTokenizer
from transformers import AutoModelForSeq2SeqLM


# <-----------------------------------Set up Streamlit app------------------------------------>
st.set_page_config(page_title="Corporate Training Knowledge Hub", layout="wide")
st.title("Corporate Training Knowledge Hub")

# <------------------------------------Initialize components------------------------------------->
# Load BART model for summarization
bart_model_name = "facebook/bart-large-cnn"
bart_tokenizer = BartTokenizer.from_pretrained(bart_model_name)
bart_model = BartForConditionalGeneration.from_pretrained(bart_model_name)



# Initialize the LLaMA model using Ollama
llm = Ollama(model="llama3.2")  # Replace with your locally installed LLaMA model


    


# Load embedding model for document processing
embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
vectorstore = None
document_store = []

# <---------------------------------------------Define tabs for functionalities------------------------------------>
tabs = st.tabs(["Upload Files", "Original Context", "Document Summarization", "Interactive Q&A", "Word Cloud", "Compare Docs"])

# <--------------------------------------------------Upload and process files------------------------------------->
def process_files(uploaded_files):
    global vectorstore
    for uploaded_file in uploaded_files:
        file_type = uploaded_file.name.split(".")[-1]
        combined_content = ""
        document = None

        with tempfile.NamedTemporaryFile(delete=False) as temp_file:
            temp_file.write(uploaded_file.read())
            temp_file_path = temp_file.name

        if file_type == "pdf":
            loader = PyPDFLoader(temp_file_path)
            all_pages = loader.load()
            combined_content = " ".join([page.page_content for page in all_pages])
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "pptx":
            presentation = Presentation(temp_file_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        combined_content += shape.text_frame.text + " "
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "txt":
            with open(temp_file_path, "r", encoding="utf-8") as file:
                combined_content = file.read()
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        elif file_type == "xlsx":
            excel_data = pd.read_excel(temp_file_path)
            combined_content = excel_data.to_string(index=False)
            document = Document(page_content=combined_content, metadata={"name": uploaded_file.name})

        if document is None:
            st.warning(f"File type '{file_type}' is not supported.")
            os.remove(temp_file_path)
            continue

        document_store.append(document)

        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_documents([document])

        if vectorstore is None:
            vectorstore = FAISS.from_documents(chunks, embedding_model)
        else:
            vectorstore.add_documents(chunks)

        os.remove(temp_file_path)

# <----------------------------------------------------Summarization function------------------------------------->
def summarize_text(text, max_length=500, min_length=30):
    """
    Summarizes the provided text using facebook/bart-large-cnn.
    """
    inputs = bart_tokenizer.encode("summarize: " + text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = bart_model.generate(inputs, max_length=max_length, min_length=min_length, length_penalty=2.0, num_beams=4, early_stopping=True)
    summary = bart_tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary


# <------------------------------------------------------Interactive Q&A Functionality----------------------------------->


def answer_question_with_llama(question):
    """
    Answers a question using the locally installed LLaMA model with Ollama.
    """
    # Guard clause for empty/None question
    if not question:
        return "Please provide a question to answer."
    
    # Guard clause for vectorstore
    if not vectorstore:
        return "No documents indexed for retrieval. Please upload files first."

    try:
        # Retrieve relevant documents
        retriever = vectorstore.as_retriever(search_type="similarity", search_kwargs={"k": 3})
        retrieved_docs = retriever.get_relevant_documents(question)

        if not retrieved_docs or len(retrieved_docs) == 0:
            return "No relevant answers found for your question."
        
        prompt = ChatPromptTemplate.from_messages([
            MessagesPlaceholder(variable_name="chat_history"),
            ("user", "{input}"),
            ("user", "Given the above conversation, generate a search query to look up in order to get information relevant to the conversation")
        ])

        history_retriever_chain = create_history_aware_retriever(llm, retriever, prompt)

        answer_prompt = ChatPromptTemplate.from_messages([
            ("system", "You are a helpful assistant. Answer the user's questions strictly based on the following context. "
             "If the answer cannot be found in the context, reply with 'I'm sorry, I don't have enough information from the provided documents to answer that.'\n\n{context}"),
            MessagesPlaceholder(variable_name="chat_history"),
            ("user", "{input}")
        ])

        # Create the document processing chain
        document_chain = create_stuff_documents_chain(llm, answer_prompt)

        # Create the final conversational retrieval chain
        conversational_retrieval_chain = create_retrieval_chain(
            history_retriever_chain, 
            document_chain
        )

        # Initialize chat history if not exists
        chat_history = []

        # Invoke the chain with the question
        response = conversational_retrieval_chain.invoke({
            'chat_history': chat_history,
            "input": question
        })

        # Update chat history
        chat_history.append((
            HumanMessage(content=question), 
            AIMessage(content=response["answer"])
        ))

        return response['answer']

    except Exception as e:
        return f"An error occurred while processing your question: {str(e)}"





# <-------------------------------------------- Word Cloud Function---------------------------------->
def generate_word_cloud(text):
    """
    Generates a word cloud visualization for the provided text.
    """
    wordcloud = WordCloud(background_color="white", width=800, height=400, max_words=200, colormap="viridis").generate(text)
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation="bilinear")
    plt.axis("off")
    st.pyplot(plt)

# <------------------------------------------------------------Compare Documents------------------------------------------------------>
def compare_documents():
    """
    Compares the first two documents in the document store for similarities and differences.
    """
    if len(document_store) < 2:
        st.error("Please upload at least two documents to compare.")
        return

    doc1_content = document_store[0].page_content
    doc2_content = document_store[1].page_content

    differ = HtmlDiff()
    html_diff = differ.make_file(doc1_content.splitlines(), doc2_content.splitlines(), 
                                 fromdesc=document_store[0].metadata['name'], todesc=document_store[1].metadata['name'])
    st.write("### Comparison Result")
    st.components.v1.html(html_diff, height=600, scrolling=True)
    similarity = SequenceMatcher(None, doc1_content, doc2_content).ratio()
    st.write(f"### Similarity Score: {similarity:.2%}")

# <-------------------------------------------------------Main App-------------------------------->
st.sidebar.header("Welcome!")
st.sidebar.info("Upload corporate training documents, explore their contents, get concise summaries, generate word clouds, and ask interactive questions!")

with tabs[0]:
    st.header("Upload Files")
    uploaded_files = st.file_uploader("Upload corporate documents (PDF, PPTX, TXT, XLSX)", 
                                      type=["pdf", "pptx", "txt", "xlsx"], accept_multiple_files=True)
    if uploaded_files:
        process_files(uploaded_files)
        st.success("Files processed successfully!")

with tabs[1]:
    st.header("Original Context")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            st.text_area(f"Content_{doc.metadata['name']}", doc.page_content, height=300)
    else:
        st.info("Please upload files to display their content.")

with tabs[2]:
    st.header("Document Summarization")
    if document_store:
        for doc in document_store:
            st.write(f"### {doc.metadata['name']}")
            summary = summarize_text(doc.page_content)
            st.write("### Summary:")
            st.write(summary)
    else:
        st.info("Please upload files to summarize.")

with tabs[3]:
    st.header("Interactive Q&A")
    
    # Check for vectorstore
    if not vectorstore:
        st.info("Please upload documents to enable the Q&A functionality.")
        st.stop()

    # Initialize session state for messages if not exists
    if 'messages' not in st.session_state:
        st.session_state.messages = []

    # Display existing messages
    for message in st.session_state.messages:
        with st.chat_message(message['role']):
            st.markdown(message['content'])

    # Handle new user input
    if prompt := st.chat_input('Ask questions about the uploaded document(s)'):
        st.session_sate.messages.append({'role': 'User', 'content': prompt})

        with st.chat_message('user'):
            st.markdown(prompt)

        with st.chat_message('bot'):
            with st.spinner('Generating response........'):
                result = st.write_stream(answer_question_with_llama())
    
        st.session_state.messages.append({'role': 'assistant', 'content': result})

with tabs[4]:
    st.header("Word Cloud")
    if document_store:
        text_data = " ".join([doc.page_content for doc in document_store])
        st.write("### Word Cloud")
        generate_word_cloud(text_data)
    else:
        st.info("Please upload files to generate a word cloud.")

with tabs[5]:
    st.header("Compare Documents")
    compare_documents()